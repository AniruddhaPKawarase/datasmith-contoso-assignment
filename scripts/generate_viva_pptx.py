"""Generate the BITS Outline Viva presentation — visual-first, 7 slides.

Output: docs/viva/Outline_Viva_2024AA05175.pptx

Design principles
-----------------
- Seven slides only (1 title + 5 content + 1 thanks).
- One idea per slide; no paragraph prose.
- Visual primitives: cards, big numbers, status dots, arrows, side-by-side
  comparisons. Tables are last resort.
- Lots of white space. The viva is verbal — slides are a backdrop, not a
  script.

Run from project root:
    python scripts/generate_viva_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── design tokens ────────────────────────────────────────────────────

NAVY = RGBColor(0x1E, 0x3A, 0x8A)
CHARCOAL = RGBColor(0x1F, 0x29, 0x37)
SLATE = RGBColor(0x47, 0x55, 0x69)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)
SOFT_NAVY = RGBColor(0xDB, 0xEA, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
TEAL = RGBColor(0x0F, 0x76, 0x6E)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.6)
TITLE_Y = Inches(0.45)
FOOTER_Y = Inches(7.1)
CONTENT_TOP = Inches(1.4)

# ── primitives ───────────────────────────────────────────────────────


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def textbox(slide, x, y, w, h, text, *, size=14, bold=False,
            color=CHARCOAL, align=PP_ALIGN.LEFT, font="Calibri",
            anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Pt(0))
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def shape(slide, kind, x, y, w, h, *, fill=None, line=None,
          line_width=0.75):
    s = slide.shapes.add_shape(kind, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_width)
    s.shadow.inherit = False
    return s


def card(slide, x, y, w, h, *, fill=WHITE, border=SOFT_NAVY,
         border_width=1.0):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
                 fill=fill, line=border, line_width=border_width)


def title_bar(slide, title: str):
    # accent dash
    shape(slide, MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(1.05),
          Inches(0.45), Inches(0.06), fill=NAVY)
    textbox(slide, MARGIN_X, TITLE_Y,
            SLIDE_W - 2 * MARGIN_X, Inches(0.55),
            title, size=28, bold=True, color=NAVY)


def footer(slide, page: int, total: int):
    textbox(slide, MARGIN_X, FOOTER_Y, Inches(9.0), Inches(0.25),
            "Aniruddha Prakash Kawarase   2024AA05175   "
            "BITS Pilani M.Tech AI/ML   Outline Viva  ·  26 May 2026",
            size=9, color=MUTED)
    textbox(slide, SLIDE_W - MARGIN_X - Inches(1.0), FOOTER_Y,
            Inches(1.0), Inches(0.25),
            f"{page} / {total}",
            size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def numbered_chip(slide, x, y, w, h, number: str, label: str):
    """One numbered objective/scope card."""
    card(slide, x, y, w, h, fill=WHITE, border=SOFT_NAVY)
    # number badge
    badge_d = Inches(0.55)
    shape(slide, MSO_SHAPE.OVAL,
          x + Inches(0.25), y + Inches(0.3), badge_d, badge_d,
          fill=NAVY)
    textbox(slide, x + Inches(0.25), y + Inches(0.3),
            badge_d, badge_d, number,
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
    # label
    textbox(slide, x + Inches(1.0), y + Inches(0.3),
            w - Inches(1.2), h - Inches(0.4), label,
            size=13, bold=True, color=CHARCOAL,
            anchor=MSO_ANCHOR.MIDDLE)


def status_dot(slide, x, y, done: bool):
    """Small filled circle indicating phase status."""
    d = Inches(0.30)
    shape(slide, MSO_SHAPE.OVAL, x, y, d, d,
          fill=GREEN if done else MUTED)


def metric_card(slide, x, y, w, h, value: str, label: str,
                *, value_color=NAVY):
    card(slide, x, y, w, h, fill=WHITE, border=SOFT_NAVY, border_width=1.25)
    textbox(slide, x, y + Inches(0.25), w, Inches(0.95),
            value, size=40, bold=True, color=value_color,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, x, y + h - Inches(0.65), w, Inches(0.45),
            label, size=11, color=SLATE,
            align=PP_ALIGN.CENTER)


def arrow_down(slide, x, y, h=Inches(0.30)):
    shape(slide, MSO_SHAPE.DOWN_ARROW, x, y, Inches(0.3), h,
          fill=NAVY)


# ── slides ───────────────────────────────────────────────────────────


def slide_title(prs):
    s = blank(prs)
    # left navy band
    shape(s, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
          Inches(0.45), SLIDE_H, fill=NAVY)
    # course header
    textbox(s, Inches(1.0), Inches(0.8), Inches(11), Inches(0.4),
            "BITS Pilani  ·  WILP  ·  Outline Viva",
            size=14, color=SLATE)
    textbox(s, Inches(1.0), Inches(1.15), Inches(11), Inches(0.4),
            "AIMLCZG628T — Dissertation",
            size=14, color=SLATE)
    # title
    textbox(s, Inches(1.0), Inches(2.2), Inches(11.6), Inches(1.0),
            "Domain-Aware Multi-Agent",
            size=48, bold=True, color=NAVY)
    textbox(s, Inches(1.0), Inches(3.05), Inches(11.6), Inches(1.0),
            "NL-to-SQL for Supply Chain",
            size=48, bold=True, color=NAVY)
    # accent rule
    shape(s, MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.15),
          Inches(2.0), Inches(0.06), fill=NAVY)
    # student
    textbox(s, Inches(1.0), Inches(4.55), Inches(11), Inches(0.5),
            "Aniruddha Prakash Kawarase",
            size=22, bold=True, color=CHARCOAL)
    textbox(s, Inches(1.0), Inches(5.05), Inches(11), Inches(0.4),
            "BITS ID  2024AA05175   ·   M.Tech AI & ML",
            size=14, color=SLATE)
    textbox(s, Inches(1.0), Inches(6.3), Inches(11), Inches(0.4),
            "May 2026",
            size=12, color=MUTED)


def slide_objectives_scope(prs, page, total):
    s = blank(prs)
    title_bar(s, "Objectives & Scope")

    # 6 objective chips in a 2x3 grid
    objs = [
        ("1", "Domain-axis decomposition"),
        ("2", "Cross-DB federation"),
        ("3", "Deterministic temporal reasoning"),
        ("4", "Source-agent self-correction"),
        ("5", "SCM-SQL benchmark (500+ queries)"),
        ("6", "Evaluation vs MAC-SQL / MARS-SQL / CHASE-SQL"),
    ]
    chip_w = Inches(4.0)
    chip_h = Inches(1.15)
    col_gap = Inches(0.15)
    row_gap = Inches(0.20)
    start_x = MARGIN_X
    start_y = Inches(1.55)
    for i, (n, label) in enumerate(objs):
        row, col = divmod(i, 3)
        x = start_x + col * (chip_w + col_gap)
        y = start_y + row * (chip_h + row_gap)
        numbered_chip(s, x, y, chip_w, chip_h, n, label)

    # scope chips at bottom
    scope_y = Inches(4.65)
    textbox(s, MARGIN_X, scope_y - Inches(0.4),
            Inches(12), Inches(0.35), "Scope",
            size=14, bold=True, color=NAVY)

    scope_items = [
        ("5", "domain agents"),
        ("6", "complexity levels"),
        ("8", "complexity dimensions"),
        ("498", "Odoo tables"),
        ("PG · DuckDB · Redis", "federated stack"),
    ]
    chip_w2 = Inches(2.45)
    chip_h2 = Inches(1.5)
    gap = Inches(0.10)
    total_w = chip_w2 * len(scope_items) + gap * (len(scope_items) - 1)
    start_x2 = (SLIDE_W - total_w) / 2
    for i, (big, small) in enumerate(scope_items):
        x = start_x2 + i * (chip_w2 + gap)
        card(s, x, scope_y, chip_w2, chip_h2, fill=LIGHT_BG, border=SOFT_NAVY)
        textbox(s, x, scope_y + Inches(0.15), chip_w2, Inches(0.85),
                big, size=26 if len(big) <= 4 else 18, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x, scope_y + chip_h2 - Inches(0.55), chip_w2, Inches(0.4),
                small, size=11, color=SLATE, align=PP_ALIGN.CENTER)

    footer(s, page, total)


def slide_architecture(prs, page, total):
    s = blank(prs)
    title_bar(s, "Methodology — Architecture")

    layers = [
        ("1.  Query Understanding",
         "Router  ·  Temporal Parser  ·  Ambiguity  ·  References"),
        ("2.  Domain Specialists",
         "Inventory  ·  Logistics  ·  Finance  ·  Demand"),
        ("3.  Composer  (sqlglot AST)",
         "CTE wrapping  +  shared-key INNER JOIN / CROSS JOIN"),
        ("4.  Compliance  (sqlglot AST)",
         "Per-SELECT RBAC injection  +  audit log"),
        ("5.  Validator",
         "Syntax  →  Postgres EXPLAIN  →  Business rules"),
    ]

    layer_w = Inches(11.5)
    layer_h = Inches(0.85)
    arrow_h = Inches(0.18)
    total_block = layer_h * len(layers) + arrow_h * (len(layers) - 1)
    start_y = CONTENT_TOP + (Inches(5.4) - total_block) / 2
    x = (SLIDE_W - layer_w) / 2

    for i, (head, sub) in enumerate(layers):
        y = start_y + i * (layer_h + arrow_h)
        card(s, x, y, layer_w, layer_h, fill=SOFT_NAVY, border=NAVY,
             border_width=1.0)
        textbox(s, x + Inches(0.25), y, Inches(4.5), layer_h,
                head, size=14, bold=True, color=NAVY,
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + Inches(4.5), y, layer_w - Inches(4.7), layer_h,
                sub, size=12, color=CHARCOAL,
                anchor=MSO_ANCHOR.MIDDLE)
        if i < len(layers) - 1:
            arrow_down(s, x + layer_w / 2 - Inches(0.15),
                       y + layer_h - Inches(0.02), h=arrow_h + Inches(0.02))

    # retry loop annotation on the right
    annotation_x = x + layer_w + Inches(0.15)
    textbox(s, annotation_x, start_y, Inches(1.5), Inches(5.0),
            "↑\nretry to\noriginating\nagent\n(max 3)",
            size=10, color=SLATE, anchor=MSO_ANCHOR.MIDDLE,
            align=PP_ALIGN.LEFT)

    footer(s, page, total)


def slide_novelty(prs, page, total):
    s = blank(prs)
    title_bar(s, "Methodology — Novelty")

    # Two columns
    col_w = Inches(5.9)
    col_h = Inches(4.6)
    gap = Inches(0.5)
    left_x = (SLIDE_W - 2 * col_w - gap) / 2
    right_x = left_x + col_w + gap
    y = Inches(1.7)

    # LEFT — existing approach
    card(s, left_x, y, col_w, col_h, fill=LIGHT_BG, border=MUTED)
    textbox(s, left_x, y + Inches(0.25), col_w, Inches(0.45),
            "Existing multi-agent NL-to-SQL",
            size=12, color=SLATE, align=PP_ALIGN.CENTER)
    textbox(s, left_x, y + Inches(0.7), col_w, Inches(0.7),
            "Decompose by SQL PIPELINE STAGE",
            size=20, bold=True, color=SLATE, align=PP_ALIGN.CENTER)
    # stage chips
    stages = ["Schema Link", "Decompose", "Generate", "Refine"]
    chip_h = Inches(0.55)
    chip_gap = Inches(0.12)
    chips_y = y + Inches(1.6)
    chip_w = col_w - Inches(1.2)
    chip_x = left_x + Inches(0.6)
    for i, stage in enumerate(stages):
        cy = chips_y + i * (chip_h + chip_gap)
        card(s, chip_x, cy, chip_w, chip_h, fill=WHITE, border=MUTED)
        textbox(s, chip_x, cy, chip_w, chip_h, stage,
                size=14, bold=True, color=SLATE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # RIGHT — this dissertation
    card(s, right_x, y, col_w, col_h, fill=SOFT_NAVY, border=NAVY,
         border_width=1.25)
    textbox(s, right_x, y + Inches(0.25), col_w, Inches(0.45),
            "This dissertation",
            size=12, color=NAVY, align=PP_ALIGN.CENTER, bold=True)
    textbox(s, right_x, y + Inches(0.7), col_w, Inches(0.7),
            "Decompose by BUSINESS DOMAIN",
            size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    domains = ["Inventory", "Logistics", "Finance", "Demand"]
    for i, d in enumerate(domains):
        cy = chips_y + i * (chip_h + chip_gap)
        card(s, chip_x + (right_x - left_x), cy, chip_w, chip_h,
             fill=WHITE, border=NAVY)
        textbox(s, chip_x + (right_x - left_x), cy, chip_w, chip_h, d,
                size=14, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # bottom note
    textbox(s, MARGIN_X, Inches(6.5),
            SLIDE_W - 2 * MARGIN_X, Inches(0.4),
            "Each agent owns its schema slice, glossary, and few-shots. "
            "Error blame routes to the originating agent — not a global retry.",
            size=11, color=SLATE, align=PP_ALIGN.CENTER)

    footer(s, page, total)


def slide_progress(prs, page, total):
    s = blank(prs)
    title_bar(s, "Progress & Key Findings")

    # ── Phase tracker (10 boxes horizontally) ──
    textbox(s, MARGIN_X, Inches(1.5), Inches(12), Inches(0.4),
            "Phase status  ·  7 of 10 complete  ·  ~5 weeks ahead of midterm",
            size=12, color=SLATE)

    n_phases = 10
    box_w = Inches(1.15)
    box_h = Inches(0.85)
    gap = Inches(0.05)
    total_w = box_w * n_phases + gap * (n_phases - 1)
    start_x = (SLIDE_W - total_w) / 2
    tracker_y = Inches(1.95)
    done_count = 7
    phase_labels = [
        "Setup", "Schema", "Framework", "Specialists", "Composer",
        "Temporal", "Multi-turn", "Benchmark", "Frontend", "Eval+Paper",
    ]
    for i in range(n_phases):
        x = start_x + i * (box_w + gap)
        done = i < done_count
        card(s, x, tracker_y, box_w, box_h,
             fill=NAVY if done else LIGHT_BG,
             border=NAVY if done else MUTED)
        textbox(s, x, tracker_y + Inches(0.08), box_w, Inches(0.3),
                f"Phase {i + 1}", size=9, bold=True,
                color=WHITE if done else SLATE,
                align=PP_ALIGN.CENTER)
        textbox(s, x, tracker_y + Inches(0.40), box_w, Inches(0.45),
                phase_labels[i],
                size=10, bold=True,
                color=WHITE if done else SLATE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # legend
    legend_y = tracker_y + box_h + Inches(0.15)
    legend_x = start_x
    status_dot(s, legend_x, legend_y + Inches(0.04), True)
    textbox(s, legend_x + Inches(0.4), legend_y, Inches(1.5), Inches(0.35),
            "complete", size=10, color=SLATE)
    status_dot(s, legend_x + Inches(2.0), legend_y + Inches(0.04), False)
    textbox(s, legend_x + Inches(2.4), legend_y, Inches(1.5), Inches(0.35),
            "pending", size=10, color=SLATE)

    # ── Key findings — 4 metric cards ──
    textbox(s, MARGIN_X, Inches(4.4), Inches(12), Inches(0.4),
            "Live evidence — measured artefacts",
            size=12, color=SLATE)

    metrics = [
        ("498", "Odoo tables introspected"),
        ("4 / 4", "queries pass live EXPLAIN"),
        ("210", "unit tests  ·  87% coverage"),
        ("$0.10", "total LLM dev spend"),
    ]
    card_w = Inches(2.85)
    card_h = Inches(1.85)
    gap_m = Inches(0.20)
    total_w_m = card_w * 4 + gap_m * 3
    start_x_m = (SLIDE_W - total_w_m) / 2
    metric_y = Inches(4.85)
    for i, (val, label) in enumerate(metrics):
        x = start_x_m + i * (card_w + gap_m)
        metric_card(s, x, metric_y, card_w, card_h, val, label)

    footer(s, page, total)


def slide_future(prs, page, total):
    s = blank(prs)
    title_bar(s, "Future Work")

    # 3 phase cards horizontally
    phases = [
        ("Phase 8", "SCM-SQL Benchmark",
         "500+ NL/SQL pairs  ·  6 complexity levels  ·  Odoo + DataCo",
         "5 Jul 2026"),
        ("Phase 9", "Next.js Frontend",
         "Chat UI  ·  SQL viewer  ·  agent timeline  ·  lineage panel",
         "16 Jul 2026"),
        ("Phase 10", "Eval + Paper + Dissertation",
         "Run vs 4 baselines  ·  ablations  ·  ACL/EMNLP/COLING draft",
         "1 Aug 2026"),
    ]
    card_w = Inches(3.95)
    card_h = Inches(4.2)
    gap = Inches(0.25)
    total_w = card_w * 3 + gap * 2
    start_x = (SLIDE_W - total_w) / 2
    y = Inches(1.85)
    for i, (phase, name, detail, deadline) in enumerate(phases):
        x = start_x + i * (card_w + gap)
        card(s, x, y, card_w, card_h, fill=WHITE, border=NAVY,
             border_width=1.25)
        # phase header band
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
              x, y, card_w, Inches(0.75),
              fill=NAVY, line=NAVY)
        textbox(s, x, y, card_w, Inches(0.75),
                phase, size=18, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x + Inches(0.25), y + Inches(1.0),
                card_w - Inches(0.5), Inches(0.7), name,
                size=18, bold=True, color=CHARCOAL,
                align=PP_ALIGN.CENTER)
        textbox(s, x + Inches(0.25), y + Inches(1.85),
                card_w - Inches(0.5), Inches(1.7), detail,
                size=12, color=SLATE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # deadline pill
        pill_w = Inches(2.0)
        pill_x = x + (card_w - pill_w) / 2
        pill_y = y + card_h - Inches(0.7)
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE,
              pill_x, pill_y, pill_w, Inches(0.45),
              fill=SOFT_NAVY, line=NAVY)
        textbox(s, pill_x, pill_y, pill_w, Inches(0.45),
                deadline, size=11, bold=True, color=NAVY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    footer(s, page, total)


def slide_thanks(prs, page, total):
    s = blank(prs)

    # giant centred thank-you
    textbox(s, Inches(0), Inches(2.4),
            SLIDE_W, Inches(1.5), "Thank You",
            size=72, bold=True, color=NAVY,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, Inches(0), Inches(3.9), SLIDE_W, Inches(0.6),
            "Questions welcome.",
            size=22, color=SLATE, align=PP_ALIGN.CENTER)

    # accent line
    shape(s, MSO_SHAPE.RECTANGLE,
          (SLIDE_W - Inches(1.5)) / 2, Inches(4.7),
          Inches(1.5), Inches(0.05), fill=NAVY)

    # compact key references at bottom
    textbox(s, MARGIN_X, Inches(5.2),
            SLIDE_W - 2 * MARGIN_X, Inches(0.3),
            "Key references",
            size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    refs = (
        "MAC-SQL (COLING 2025)   ·   MARS-SQL (2025)   ·   "
        "CHASE-SQL (ICLR 2025)   ·   Spider 2.0 (ICLR 2025)   ·   "
        "AmbiSQL (2025)   ·   NL2SQL Survey (VLDB 2025)"
    )
    textbox(s, MARGIN_X, Inches(5.55),
            SLIDE_W - 2 * MARGIN_X, Inches(0.4),
            refs, size=10, color=SLATE, align=PP_ALIGN.CENTER)

    footer(s, page, total)


# ── main ─────────────────────────────────────────────────────────────


def main() -> int:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 7
    slide_title(prs)
    slide_objectives_scope(prs, 2, total)
    slide_architecture(prs, 3, total)
    slide_novelty(prs, 4, total)
    slide_progress(prs, 5, total)
    slide_future(prs, 6, total)
    slide_thanks(prs, 7, total)

    out_dir = Path(__file__).resolve().parent.parent / "docs" / "viva"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "Outline_Viva_2024AA05175.pptx"
    prs.save(out_path)

    size_kb = out_path.stat().st_size / 1024
    print(f"Wrote {out_path.relative_to(out_dir.parent.parent)}  "
          f"({size_kb:.1f} KB, {total} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
